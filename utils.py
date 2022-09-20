#!/usr/bin/env python
# -*- coding: utf-8 -*-

__mtime__ = '2019/9/8'

# from gevent import monkey
# monkey.patch_all()
# import gevent

import os,shutil,zipfile,rarfile,random
import time
from threading import Thread
from multiprocessing import Process,Pool
from datetime import datetime

import win32com
import win32gui
import win32con
import win32com.client


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import glob
from docx import Document


import config as cfg
import common as com


# 压缩一个目录
def zip_dir(dirname,zipfilename,mode):
    filelist = []
    if os.path.isfile(dirname):
        filelist.append(dirname)
    else :
        for root, dirs, files in os.walk(dirname):
            for dir in dirs:
                filelist.append(os.path.join(root,dir))
            for name in files:
                filelist.append(os.path.join(root, name))

    zf = zipfile.ZipFile(zipfilename, mode, zipfile.zlib.DEFLATED)
    for tar in filelist:
        # arcname = tar[len(dirname):]
        # arcname=dirname
        _,arcname=os.path.split(dirname)

        #print arcname
        zf.write(tar,arcname)
    zf.close()

def zip_dir_w(dirname,zipfilename):
    return zip_dir(dirname,zipfilename,"w")

def zip_dir_a(dirname,zipfilename):
    return zip_dir(dirname,zipfilename,"a")



# 解压后的文件夹与原来的zip文件同名且在相同目录下
# 确保windows下解压后的文件夹名不乱码
file_encoding = "utf-8"
# if platform.system() == "Windows":
#     file_encoding = "gbk"

# 将zip文件解压到其所在目录
def unzip(file):

    file_name, ext = os.path.splitext(file)
    try:
        if ext == ".zip":
            print ('unzip', file)
        f = zipfile.ZipFile(file)
        # f.extractall(path=file_name.encode(file_encoding)) # 通过path指定解压的路径
        f.extractall(path=r'{}'.format(file_name)) # 通过path指定解压的路径
    except Exception as e:
        print(e)


def un_rar(file_name):
    """unrar zip file"""
    try:
        rar = rarfile.RarFile(file_name)
        dir_name=file_name.replace('.rar','')
        if os.path.isdir(dir_name):
            pass
        else:
            os.mkdir(dir_name)
        rar.extractall(path=dir_name)
        rar.close()
    except Exception as e:
        print(e)

# 递归删除目录及目录下文件
def removeDir(dirPath):
    if not os.path.isdir(dirPath):
       return
    files = os.listdir(dirPath)
    try:
       for file in files:
           filePath = os.path.join(dirPath,file)
           if os.path.isfile(filePath):
               os.remove(filePath)
           elif os.path.isdir(filePath):
                removeDir(filePath)
       os.rmdir(dirPath)
    except Exception as e:
       print(e)


def drop_ppt_filter_page(ppt_path,only_last=False):
    # 读取ppt
    try:
        prs = Presentation(ppt_path)

        # 查看一共几页
        slides = prs.slides
        number_pages = len(slides)
        print(number_pages)


        # 删除最后一、二页
        if isinstance(only_last,bool):
            if only_last==True:
                rId = prs.slides._sldIdLst[-1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[-1]
            else:
                for i in range(2):
                    rId = prs.slides._sldIdLst[-1].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[-1]
        else:
            if isinstance(only_last,int):
                rId = prs.slides._sldIdLst[only_last].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[only_last]


        # 保存新的ppt
        prs.save(ppt_path)
    except Exception as e:
        print(e)


def many_drop_ppt_filter_page(path=cfg.save_path,file_type=('.pptx','.ppt'),only_last=True):
    for root,dirs,files in os.walk(path):
        for f in files:
            if f.endswith(file_type):
                full_path=root+'/'+f
                drop_ppt_filter_page(full_path,only_last=only_last)


def filter_ppt_text(path,replace_from):
    li=[]
    for root,dirs,files in os.walk(path):
        for f in files:
            # handle_ppt_text(f,path,replace_from)
            p=Process(target=handle_ppt_text,args=(f,path,replace_from,))
            li.append(p)
            p.start()
        #     t=Thread(target=handle_ppt_text,args=(f,path,replace_from,))
        #     li.append(t)
        #     t.start()
        #
        for n in li:
            n.join()
    # gevent.joinall([
    #     gevent.spawn(handle_ppt_text,f,path,replace_from) for root,dirs,files in os.walk(path) for f in files
    # ])


    print('过滤完成')


def handle_ppt_text(f,path,replace_from):
    if f.endswith(('.pptx','.ppt')):
                full_path=path+'/'+f
                print(full_path)
                try:
                    prs = Presentation(full_path)

                    for j in range(0,len(prs.slides)):
                        # 去备注
                        note_slide=prs.slides[j].notes_slide
                        note_slide_handle(note_slide,f)

                        #PPT文本过滤与替换
                        shape=prs.slides[j].shapes
                        shape_handle(shape,f,replace_from)
                        upadate_date(shape,f,'202X')

                    prs.save(full_path)
                except Exception as e:
                    print(e)



def upadate_date(shape,file_name,date='202X'):
    li=['201{}'.format(i) for i in range(0,9)]
    li.append('2020')
    # li.append('2021')
    shape_handle(shape,file_name,li,replace_to=date)



def replace_text(paragraph,run_frame,replace_from,file_name,replace_to=None):
    if isinstance(replace_from,list):
        for replace_from_one in replace_from:
            replace_text(paragraph,run_frame,replace_from_one,file_name,replace_to)

    else:
        if replace_from in run_frame.text:
            if replace_to==None:
                replace_to=com.choice_text_to_replace(replace_from)
                run_frame.text=run_frame.text.replace(replace_from,replace_to)

                print('[{}]{}:已去出处信息'.format(datetime.now().strftime('%Y-%m-%d %H:%M:%S'),file_name))
            else:
                run_frame.text=run_frame.text.replace(replace_from,replace_to)
                print('[{}]{}:已进行自定义替换{}'.format(datetime.now().strftime('%Y-%m-%d %H:%M:%S'),file_name,replace_to))


def change_shape_text(shape,replace_from,file_name,replace_to=None):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shape_in_group in shape.shapes:
            change_shape_text(shape_in_group,replace_from,file_name,replace_to)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for cell in shape.table.iter_cells():
            text_frame = cell.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    replace_text(run,text_frame,replace_from,file_name,replace_to)
    else:
        if shape.has_text_frame:
            shape_text_frame = shape.text_frame
            for paragraph in shape_text_frame.paragraphs:
                for run in paragraph.runs:
                    replace_text(paragraph,run,replace_from,file_name,replace_to)



def shape_handle(shapes,f,replace_from,replace_to=None):
        for i in range(0,len(shapes)):
            change_shape_text(shapes[i],replace_from,f,replace_to)




def ppt_to_jpg(path):
    try:
        # from_path='./download_file/测试用ppt.pptx'
        wpp = win32com.client.Dispatch("Kwpp.Application")
        # wpp = win32com.client.Dispatch("PowerPoint.Application")
        # wpp.Visible=0
        # wpp.DisplayAlerts=0
        for root,dirs,files in os.walk(path):
            for file in files:
                if file.endswith(('.ppt','.pptx')):
                    dir_name,_=os.path.splitext(file)
                    changed_dirs=com.get_changed(cfg.pic_path,'dirs')
                    if dir_name in changed_dirs:
                        continue

                    from_path=os.path.join(path,file)

                    ppt = wpp.Presentations.Open(from_path)
                    to_path=from_path.replace('.pptx','.jpg')
                    to_path=to_path.replace('.ppt','.jpg')
                    file_name=os.path.split(to_path)[1]

                    to_path='{}/{}'.format(cfg.pic_path,file_name)

                    ppt.SaveAs(to_path,17)

                print('{}:图片转换完成'.format(file))
    except Exception as e:
        print(e)
    finally:
        # ppt.Close()
        wpp.Quit()

def remove_empty_dirs(path):
    for root,dirs,files in os.walk(path):
        for d in dirs:
            full_path=root+'/'+d

            d1=os.listdir(full_path)
            if d1==[]:
                removeDir(full_path)
                print("{}:是空文件夹,已删除".format(d))
    print('所有空文件夹已删除')


def get_rar_path(dir_path):
    dir_path=dir_path.replace('\u202a','')
    dir_li = os.listdir(dir_path)
    li=[]
    for i in dir_li:
        if i.endswith('.rar'):
            a=dir_path+i
            a=a.replace('/','\\')
            # print(a)
            li.append(a)
    print(li)
    return li

def win_handle(path,title=u'打开'):
     # win32gui

    dialog = win32gui.FindWindow('#32770', title)  # 对话框
    ComboBoxEx32 = win32gui.FindWindowEx(dialog, 0, 'ComboBoxEx32', None)
    ComboBox = win32gui.FindWindowEx(ComboBoxEx32, 0, 'ComboBox', None)
    Edit = win32gui.FindWindowEx(ComboBox, 0, 'Edit', None)  # 上面三句依次寻找对象，直到找到输入框Edit对象的句柄
    button = win32gui.FindWindowEx(dialog, 0, 'Button', None)  # 确定按钮Button

    # time.sleep(0.5)

    win32gui.SendMessage(Edit, win32con.WM_SETTEXT, None,path)  # 往输入框输入绝对地址
    win32gui.SendMessage(dialog, win32con.WM_COMMAND, 1, button)  # 按button


def get_file_path(dir_name,end_with=(".ppt",'.pptx'),to_win32=False):
    li=[]
    for root,dirs,files in os.walk(dir_name):
        for file_name in files:
            if file_name.endswith(end_with):
                full_path=os.path.join(root,file_name)
                li.append(full_path)
    if to_win32==True:
        li2=[]
        for i in li:
            i=i.replace('/','\\')
            li2.append(i)
        return li2
    else:
        return li


def filter_pic():
    a=get_file_path(cfg.pic_path,end_with=('.jpg'))
    filter_tuple=tuple(["幻灯片{}.jpg".format(i) for i in range(1,6)])
    for i in a:
        if not i.endswith(filter_tuple):
            os.remove(i)

    print('多余图片过滤完成')

def rename_file(path,add_name='random'):
    for root,dirs,files in os.walk(path):
        for f in files:
            full_name=os.path.join(root,f)
            dst_name,ext=os.path.splitext(full_name)

            if add_name=='random':
                dst_name=dst_name+str(random.randint(1,1000))
            else:
                dst_name=dst_name+str(add_name)
                # dst_name=dst_name.replace(str(add_name),'')
            final_name=dst_name+ext
            os.rename(full_name,final_name)
            print('{}:重命名成功'.format(dst_name))


def handle_file(file_path,save_path,end_with=('.doc','.docx','DOC')):
    _,ext=os.path.splitext(file_path)
    if '.zip'==ext and zipfile.is_zipfile(file_path):
        unzip(r'{}'.format(file_path))
    else:

        un_rar(r'{}'.format(file_path))
    dir_name=os.path.split(file_path)[1].replace('.zip','')
    dir_name=dir_name.replace('.rar','')

    dir_path=os.path.join(save_path,dir_name)

    for root,dirs,files in os.walk(dir_path):

        for f in files:
            if f.endswith(end_with):
                old_path=root+'/'+f

                old_name,ext=os.path.splitext(f)
                new_name=dir_name+ext
                new_path=root+'/'+new_name

                if old_name!=new_name:
                    os.rename(old_path,new_path)


                final_path=os.path.join(save_path,new_name)

                shutil.move(new_path,final_path)
                os.remove(file_path)
                removeDir(dir_path)



def handle_compress(save_path,end_with):
    for root,dirs,files in os.walk(save_path):
        for f in files:
            try:
                if f.endswith(('.rar','.zip')):
                    full_path=root+'/'+f
                    handle_file(full_path,save_path,end_with)
            except Exception as e:
                print(e)
                continue



def rename_ppt_note(ppt_path,rename=''):
    for root,dirs,files in os.walk(ppt_path):
        for f in files:
            if f.endswith(('.pptx','ppt')):
                full_path=root+'/'+f
                print(full_path)
                prs=Presentation(full_path)

                slides=prs.slides
                for s in slides:
                    note_slide=s.notes_slide
                    print(dir(note_slide))
                    note_tf=note_slide.notes_text_frame
                    if note_tf.text!='':
                       note_tf.text=note_tf.text.replace(note_tf.text,'')
                print('{}:备注已清除'.format(f))
                prs.save(full_path)


def note_slide_handle(note_slide,file_name):
    note_tf=note_slide.notes_text_frame
    if note_tf.text!='':
        note_tf.text=note_tf.text.replace(note_tf.text,'')
        print('{}:备注已清除'.format(file_name))


def mkdir(path):
    if not os.path.exists(path):
        os.mkdir(path)
        print(path,'：创建成功')
    else:
        print(path,':已存在，不能重复创建')

def auto_mkdir():
    mkdir(path=cfg.save_path)
    mkdir(path=cfg.pic_path)



def time_sleep_calc(path):
    if os.path.getsize(path)>=1500000000:
        t=600
    elif os.path.getsize(path)>=1000000000:
        t=400
    elif os.path.getsize(path)>=500000000:
        t=200
    elif os.path.getsize(path)>=200000000:
        t=60
    elif os.path.getsize(path)>=100000000:
        t=40
    elif os.path.getsize(path)>=40000000:
        t=20
    elif os.path.getsize(path)>=20000000:
        t=15
    elif os.path.getsize(path)>=5000000:
        t=10
    else:
         # 小于5M等待6秒
        t=6
    return t



if __name__ == "__main__":

    rename_file(cfg.ppt_test_path)

    # remove_empty_dirs(cfg.pic_path)

